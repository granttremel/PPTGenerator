# -*- coding: utf-8 -*-
"""
Created on Mon Jan  8 10:39:25 2024

@author: peyton2
"""

from pptx import Presentation
from pptx.shapes.picture import Picture
from pptx.util import Length, Inches, Pt
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.dml.line import LineFormat
from datetime import datetime,timedelta
from enum import Enum,IntEnum

from . import util
# from glib import gtools as gt
# from glib import improc as ip
# from glib import plotting as gplt
from glib import cosmx_strings
# from glib import logger

import traceback
import matplotlib as mpl
import matplotlib.pyplot as plt
import io
import numpy as np
import os
from os import path

#%%

class pptgenerator:
    
    
    cmkeys = ['top','bottom','left','right']
    
    def __init__(self,fp,cstr,params, contrast_data):
        
        if isinstance(fp,util.fp):
            self.fp  = fp
        else:
            self.fp = util.fp(fp)
        self.template_path = self.fp.ftot
        self.cstr = cstr
        self.nims = len(cstr)
        self.gridshape = tuple()
        self.params = params
        self.titles = []
        self.collabels = []
        self.rowlabels = []
        
        self.tmpdir = str(fp).replace('.','') + 'tmp'
        self.check_local = False
        
        self.groupby = tuple()
        
        self.contrast_data = contrast_data
        
        if params.contrast_broadcast_mode == 0:
            self.contrast_group = tuple()
        elif params.contrast_broadcast_mode == 1:
            self.contrast_group = ('W',)
        elif params.contrast_broadcast_mode == 2:
            self.contrast_group = ('F',)
        else:
            self.contrast_group = params.contrast_group
        
        self.layout_codes()
        self.validate_gridshape()
        
        self.collabels = []
        self.rowlabels = []
    
        logfp = util.fp(str(self.fp))
        logfp.fpath += '_log'
        logfp.fext = '.txt'
        # self.logger = logger.logger(str(logfp))
        # self.logger.information('initialized logger')
        # print = self.logger.print
    
    def generate(self, nstop= -1):
        
        self.layout_codes()
        self.validate_gridshape()
        
        print('Calculating contrast bounds')
        # if isinstance(self.contrast_data,list):
        self.get_contrasts()
        
        layouts = get_nstg_layouts(self.template_path,key = styles.RUO.name)
        
        self.prs = make_nstg_presentation(layouts, style = styles.RUO)
        
        if self.params.dark_mode:
            make_dark_mode(self.prs)
        
        
        print('Added title slide')
        add_nstg_title_slide(self.prs,color = 'white', title_text = self.cstr.expdir, name_text = 'automated with glib')
        
        print('Generating slides')
        try:
            self.generate_slides(nstop = nstop)
        except Exception as e:
            print(f'Exception encountered while generating slides: {str(e)}')
            traceback.print_exc()
            
        if self.params.dark_mode:
            make_dark_mode(self.prs)
            
        print(f'Saving presentation to {str(self.fp)}')
        self.save()
    
    def generate_first(self):
        self.generate(nstop = 1)
    
    def save(self, fpath = ''):
        if not fpath:
            fpath = str(self.fp)
            
        self.prs.save(fpath)
    
    def generate_slides(self, nstop = -1):
        
        nrounds = 0
        
        for k in self.codelist.keys():
            
            # print(f'NROUNDS={nrounds},N={nstop}')
            if nrounds == nstop:
                break
            
            figs = dict()
            
            for code in self.codelist[k]:
                
                selftup = code.to_tuple(*self.cstr.groups)
                
                im = self.retrieve_image(code)
                
                if isinstance(im,type(None)):
                    continue
                
                inset_ttl = code.to_string(*self.cstr.groups)
                
                ctup = code.to_tuple(*self.contrast_group)
                
                if self.params.contrast_broadcast_mode == 0:
                    vrange = np.percentile(im,self.contrast_data[0])
                    ckwarg = {'vrange':vrange}
                elif ctup in self.contrasts:
                    ckwarg = {'vrange':self.contrasts[ctup]}
                else:
                    print(ctup)
                    print(self.contrasts)
                    print('idk!!')
                
                try:
                    f = util.imshow(im,
                                write_contrast = True,
                                show = False,
                                channel = code.W,
                                title = inset_ttl,
                                **ckwarg)
                    figs[selftup] = f
                except IndexError as e:
                    print(f"Error setting contrasts. Try running set_contrast_data")
                except Exception as e:
                    print(f"Exception encountered while generating figure: {str(e)}")
                    traceback.print_exc()
                    continue
                
                if self.params.save_pngs:
                    pngpath = self.write_fig_to_png(code,f)
                    figs[selftup] = pngpath
                
            figlist = [figs[k] if k in figs else None for k in self.fig_layout]
            
            # if self.params.save_pngs:
            #     for fig in figlist:
            #         fig.t
            
            ncs_title = code.to_string_except(*self.cstr.groups)
            if self.params.max_per_slide > 0:
                nslidesperfigset = int(np.ceil(len(figs) / self.params.max_per_slide))
            else:
                nslidesperfigset = 1
            
            for n in range(nslidesperfigset):
                nmin = n*self.params.max_per_slide
                nmax = nmin + self.params.max_per_slide
                croppedcodes = self.codelist[k][nmin:nmax]
                croppedfigs = figlist[nmin:nmax]
                self.append_tiled_slide(croppedfigs,croppedcodes,ncs_title)
                
            nrounds +=1

    
    def append_tiled_slide(self, figlist,codelist, title):
        # ncs_title = code.to_string_except(*self.cstr.groups)
        ncs_title = title
        auto_title = self.get_auto_title(codelist)
        # print('title:', auto_title)
        
        if not any(figlist):
            print(f'No valid images identified for {ncs_title}. Skipping...')
            return
        # rows = len(figlist) // self.params.tile_cols
        
        margins = dict(zip(['top','left','bottom','right'],self.params.margins))
        
        ncs = add_nstg_nocontent_slide(self.prs,title_text = '')
        # ncs.shapes[0].text = ncs_title
        ncs.shapes[0].text = auto_title
        
        if self.params.auto_labels:
            collabels,rowlabels = self.get_auto_labels(codelist)
        else:
            collabels,rowlabels = self.collabels,self.rowlabels
        
        try:
            add_tiled_ims(self.prs,
                          ncs,
                          figlist,
                          ncs_title,
                          self.gridshape,
                          rowlabels = rowlabels,
                          collabels = collabels,
                          sep = self.params.sep,
                          custmargins = margins)
        except AssertionError:
            print(f'Encountered Error, cols = {self.gridshape[0]}, rows = {self.gridshape[1]}, nfigs = {len(figlist)}')
            print(len(figlist))
            print(figlist)
        except Exception as e:
            print(f'Encountered error while generating slides: {str(e)}')
            raise
    
    def estimate_run_time(self):
        runtime = self.nims * 0.2
        td = timedelta(seconds = runtime)
        print(f'Estimated run time: {str(td)}')
        return runtime
    
    def layout_codes(self):
        codelist = dict()
        fig_layout = list()
        contrasts = dict()
        
        groups = self.cstr.groups
        
        for code in self.cstr:
            redtup = code.to_tuple_except(*groups)
            selftup = code.to_tuple(*groups)
            cgrouptup = code.to_tuple(*self.contrast_group)
            
            if not redtup in codelist:
                codelist[redtup] = []
            
            codelist[redtup].append(code)
            
            if not selftup in fig_layout:
                fig_layout.append(selftup)
            
            contrasts[cgrouptup] = [] # list of vranges to be modified later by calc mode
            
        self.codelist = codelist
        self.fig_layout = fig_layout
        self.contrasts = contrasts

        self.nfigsperslide = len(self.fig_layout)
        self.nslides = len(self.codelist)

        self.cstr.reset()
    
    def set_groups(self,group1,group2):
        self.groupby = (group1,group2)
    
    def set_labels(self,collabels = [],rowlabels = []):
        self.rowlabels = rowlabels
        self.collabels = collabels
    
    def set_contrast_data(self, contrast_data):
        self.contrast_data = contrast_data
    
    def validate_contrast_data(self):
        #Validate user input
        cdlen = len(self.contrast_data)
        reqlen = len(self.contrasts)
        
        if cdlen == reqlen and isinstance(self.contrast_data,dict):
            return
        
        
        if not self.contrast_data:
            self.contrast_data = [(2,98)] * reqlen
        elif cdlen > reqlen:
            self.contrast_data = self.contrast_data[:reqlen]
        else:
            ratio = reqlen / cdlen
            if ratio == int(ratio) and ratio > 1:
                self.contrast_data *= int(ratio)
                print(f'Broadcasting input contrast data (n = {cdlen}) to required length (n = {reqlen})')
        
        assert len(self.contrast_data) == len(self.contrasts)
        
        # if not isinstance(self.contrast_data,dict):
        self.contrast_data = dict(zip(self.contrasts.keys(),self.contrast_data))
        
    def validate_gridshape(self):
        
        if self.params.max_per_slide > 0:
            rows = int(np.ceil(self.params.max_per_slide / self.params.tile_cols))
        elif self.params.max_rows > 0:
            rows = self.params.max_rows
            self.params.max_per_slide  = rows * self.params.tile_cols
        else:
            first = self.codelist[next(iter(self.codelist.keys()))]
            rows = int(np.ceil( len(first) / self.params.tile_cols))
            
        self.gridshape = (self.params.tile_cols,rows)
    
    def get_auto_title(self,codelist):
        titleid = set(cosmx_strings.cosmx_code.keys)
        common = cosmx_strings.cosmx_code.intersection(*codelist)
        titleid.intersection(common)
        return codelist[0].to_string(*titleid)
        
    def get_auto_labels(self, codelist):
        
        first = codelist

        #do cols
        for rowind in range(self.gridshape[1]):
            rowcodes = first[self.gridshape[0] * rowind : self.gridshape[0] * (rowind+1)]

            if not rowcodes:
                continue

            rowcommon = cosmx_strings.cosmx_code.intersection(*rowcodes)
            _rowid = rowcommon.intersection(set(self.cstr.groups))

        for colind in range(self.gridshape[1]):
            colcodes = first[colind::self.gridshape[0]]

            colcommon = cosmx_strings.cosmx_code.intersection(*colcodes)
            _colid = colcommon.intersection(set(self.cstr.groups))

        _colid = _colid.difference(_rowid)
        
        if len(_colid) == 0:
            collabels = []
        else:
            collabels = ['\n'.join([f'{eachcolid}={getattr(code,eachcolid)}' 
                         for eachcolid in _colid])
                         for code in first[:self.gridshape[0]]]
        
        if len(_rowid) ==0:
            rowlabels = []
        else:
            rowlabels = ['\n'.join([f'{eachrowid}={getattr(code,eachrowid)}' 
                         for eachrowid in _rowid])
                         for code in first[::self.gridshape[0]]]


        return collabels,rowlabels        
    
    def get_contrasts(self):
        if self.params.contrast_broadcast_mode == 0:
            return
        
        self.validate_contrast_data()

        if self.params.contrast_broadcast_mode == 0:
            #calculate contrasts on an individual basis during image collection
            return
        else:
            match self.params.contrast_calculate_mode:
                case 0:
                    self.get_first_contrasts()
                    self.contrasts = {k:self.contrasts[k][0] for k in self.contrasts.keys()}
                    
                case 1:
                    self.get_all_contrasts()
                    self.contrasts = {k:[np.min([ct[0] for ct in self.contrasts[k]]),np.max([ct[1] for ct in self.contrasts[k]])] for k in self.contrasts.keys()}
                
                case 2:
                    #this one can fail easily
                    self.get_all_contrasts()
                    self.contrasts = {k:[np.max([ct[0] for ct in self.contrasts[k]]),np.min([ct[1] for ct in self.contrasts[k]])] for k in self.contrasts.keys()}

                case 3:
                    self.get_all_contrasts()
                    self.contrasts = {k:[np.mean([ct[0] for ct in self.contrasts[k]]),np.mean([ct[1] for ct in self.contrasts[k]])] for k in self.contrasts.keys()}
                
            
    def get_first_contrasts(self):
        if not os.path.exists(self.tmpdir):
            os.makedirs(self.tmpdir)
        
        print('getting first contrasts')
        for code in self.cstr:
            ctup = code.to_tuple(*self.contrast_group)
            if len(self.contrasts[ctup]) == 0:
                
                im = self.retrieve_image(code, allow_local = True)

                vrange = np.percentile(im,self.contrast_data[ctup])
                self.contrasts[ctup].append(vrange)
            
        self.cstr.reset()
    
    def get_all_contrasts(self):
        if not os.path.exists(self.tmpdir):
            os.makedirs(self.tmpdir)
        
        for code in self.cstr:
            ctup = code.to_tuple(*self.contrast_group)
            
            im = self.retrieve_image(code, allow_local = True)
            
            if isinstance(im,type(None)):
                continue
            
            vrange = np.percentile(im,self.contrast_data[ctup])
            self.contrasts[ctup].append(vrange)
        
        self.cstr.reset()
        
    def write_fig_to_png(self,code,fig):
        fname = self.cstr.code_to_fname(code, W = code.W)
        fp = util.fp(os.path.join(self.tmpdir,fname))
        fp.fext = '.png'
        
        fig.savefig(fp.ftot, bbox_inches = 'tight', pad_inches = 0)
        print(f'Wrote figure with code {str(code)} to local file {fp.ftot}')
        
        return fp.ftot
    
    def retrieve_image(self, code, allow_local = True):
        if  allow_local and os.path.exists(self.tmpdir):
            fname = self.cstr.code_to_fname(code, W = code.W)
            fpath = os.path.join(self.tmpdir,fname)
            
            if os.path.exists(fpath):
                s,outim = util.read_im(fpath)
                print(f'Image {fname} retrieved from local storage')
                return outim
            
        fpath = self.cstr.code_to_fpath(code)
        
        if not os.path.exists(fpath):
            print(f'Could not open image at {fpath}. Skipping...')
            return
        
        ims = util.read_multichan_substack(fpath, keys = [code.W])
        
        if not code.W in ims:
            print(f'Failed to read channel {code.W} from {fpath}. Skipping...')
            return
        
        if self.params.crop == [(0,1),(0,1)]:
            im = ims[code.W]
        else:
            im = util.crop_bds(ims[code.W], self.params.crop, relative = True)
        
        print(f'Read image at {fpath} with W = {code.W}')
        
        fp_rem = util.fp(self.cstr.code_to_fpath(code))
        fp_loc = util.fp.from_parts(self.tmpdir,fp_rem.fname + f'_W{code.W}', fp_rem.fext)
        
        util.write_im(fp_loc,im)
        
        print(f'Wrote image to local storage at {str(fp_loc)}')
        
        return im
            
    
    def remove_temp_dir(self):
        
        self.check_local = False
        
        if not os.path.exists(self.tmpdir):
            return
        
        files = os.listdir(self.tmpdir)
        for file in files:
            fpath = os.path.join(self.tmpdir,file)
            os.remove(fpath)
            
        os.rmdir(self.tmpdir)
    
    def set_titles(self, titles):
        self.titles = titles
    
    def set_margins(self,margins):
        #[top,bottom,left,right] in inches
        self.margins = margins
    
    def set_template_path(self, template_path):
        self.template_path = template_path
    
    class constrast_broadcast_mode(IntEnum):
        SELF = 0
        BY_CHAN = 1
        BY_FOV = 2
        
    class contrast_calculate_mode(IntEnum):
        FIRST = 0 #just use the first image loaded (fast!)
        MINMIN_MAXMAX = 1 #maximum of maxima, minimum of minima
        MAXMIN_MINMAX = 2 #maximum of minima, minimum of maxima
        AVG = 3 #average across images in broadcast category
        

class pptgeneratorparams:
    
    
    defaults = {
        'margins' : [1.0,0.5,0.5,0.5], #top left bottom right
        'sep' : 0.1,
        'crop' : [(0,1),(0,1)],
        'tile_cols' : 3,
        'auto_labels' : False,
        'contrast_broadcast_mode' : 1,
        'contrast_group' : tuple(),
        'contrast_calculate_mode' : 0,
        'max_rows' : -1,
        'max_per_slide' : -1,
        'dark_mode': False,
        'save_pngs' : False
        }
    
    def __init__(self,**kwargs):
        """
        Parameter object which determines many of the behaviors of pptgenerator
        Valid kwargs:
        margins: list of double in order [top,left,bottom,right]
            sets margins of slides
        sep: double
            separation between adjacent images or shapes in pptx
        crop: list of two tuples, each with two doubles, in order [(xmin,xmax),(ymin,ymax)]
            determines relative bounds at which to crop image before inserting in ppt
        tile_cols: integer
            the number of columns in the image grid
        auto_labels: bool
            Whether to use auto labels or not. if set to false, custom labels can be provided
            or omitted entirely
        conrast_broadcast_mode
        contrast_group
        contrast_calculate_mode
        max_per_slide: integer or -1
            maximum number of images to insert in a slide
        
        

        """
        
        for k in self.defaults.keys():
            if k in kwargs:
                setattr(self,k,kwargs.get(k))
            else:
                setattr(self,k,self.defaults[k])
        
    @classmethod
    def default(cls):
        return cls(**cls.defaults.copy())

    def __str__(self):
        return str(self.__dict__)
    
class styles(Enum):
    RUO = 0
    nonRUO = 1
    IUO = 2
    

#%% interface with pptx
"""
this got out of hand!! needs to be reworked.
"""


name_formats = r'NSTG_CorporateTemplate_16x9_{}_2019.pptx'

keys = ['IOU','RUO','nonRUO']
default_name = 'Grant Tremel'
default_margins = {'top':Inches(1.25),
                   'left':Inches(0.5),
                   'bottom':Inches(0.5),
                   'right':Inches(0.5)}

_nstg_slide_layouts = ['title_white',
                      'title_green',
                      'title_blue',
                      'content',
                      'content1',
                      'columns',
                      'nocontent',
                      'columns1',
                      'content2',
                      'content3',
                      'columns2',
                      'blank',
                      'columns3',
                      'transition_blue',
                      'transition_green',
                      'transition_white']
nstg_slide_layouts = {k:i for i,k in enumerate(_nstg_slide_layouts)}

def get_nstg_layouts(template_path = '', key = ''):
    if not template_path:
        template_path = r"C:\Users\peyton2\Documents\GJT\Python\ppt templates"
        
    if key:
        name = name_formats.format(key)
        fullpath = os.path.join(template_path,name)
        if os.path.exists(fullpath):
            return fullpath
        else:
            raise KeyError(f'NSTG template with key {key} at path {template_path} with name {name} not found')
    
    outlist = list()
    for key in keys:
        name = name_formats.format(key)
        fullpath = os.path.join(template_path,name + '.potx')
        if os.path.exists(fullpath):
            outlist.append(fullpath)
        else:
            print(f'powerpoint with key {key} not found. continuing...')
    
    return outlist


def get_nstg_slide_layout(ind):
    pass

def make_nstg_presentation(templatepath, style = styles.RUO):
    # assert style in keys
    fullpath = str(templatepath)
    prs = Presentation(fullpath)
    
    # if len(prs.slides) > 0:
    #     for i in range(len(prs.slides)):
    #         prs.slides[i].remove()
    
    return prs

def generate_empty_ppt(prs):
    newprs = Presentation()
    
    cpnames = dir(prs.core_properties)
    
    for cpname in cpnames:
        att = getattr(prs.core_properties,cpname)
        
        if not att:
            continue
        
        if not callable(att):
            setattr(newprs.core_properties,cpname,att)
    
    newprs.slide_height = prs.slide_height
    newprs.slide_width = prs.slide_width
    
    # newprs.slide_layouts = prs.slide_layouts
    # newprs.slide_masters = prs.slide_masters
    
    return newprs


def make_dark_mode(prs):
    # for slide in prs.slides:
    #     slide.follow_master_background = True
    
    # for sm in prs.slide_masters:
    sm = prs.slide_master
    fill = sm.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(43,43,43)
    
    
    for i,slide in enumerate(prs.slides):
        if i == 0:
            continue
        for ph in slide.placeholders:
            # print(dir(ph))
            # print(ph.has_text_frame)
            if ph.has_text_frame:
                tf = ph.text_frame
                for pg in tf.paragraphs:
                    # print(dir(pg))
                    font = pg.font
                    font.color.rgb = RGBColor(212,212,212)
            
        for shape in slide.shapes:
            # print(shape.has_text_frame)
            if shape.has_text_frame:
                tf = shape.text_frame
                for pg in tf.paragraphs:
                    font = pg.font
                    font.color.rgb = RGBColor(212,212,212)
                # print(shape.text)

def add_nstg_title_slide(prs,color = 'white',title_text = '', name_text = '', date_text = ''):
    # assert color in 
    titlename = f'title_{color}'
    assert titlename in nstg_slide_layouts
    titleslide = prs.slides.add_slide(prs.slide_layouts[nstg_slide_layouts[titlename]])
    write_nstg_title_slide(titleslide, title_text,name_text, date_text)
    return titleslide

def add_nstg_title_slide_from_ref(prs,refprs,color = 'white',title_text = '', name_text = '', date_text = ''):
    # assert color in 
    titlename = f'title_{color}'
    assert titlename in nstg_slide_layouts
    titleslide = prs.slides.add_slide(refprs.slide_layouts[nstg_slide_layouts[titlename]])
    write_nstg_title_slide(titleslide, title_text,name_text, date_text)
    return titleslide
    
def write_nstg_title_slide(slide,title_text = '', name_text = '', date_text = ''):
    if not name_text:
        name_text = default_name
    
    if not date_text:
        date_text = datetime.now().strftime("%b %d, %Y")
        
    try:
        slide.shapes[0].text = name_text
        slide.shapes[1].text = date_text
        slide.shapes[2].text = title_text
    except IndexError:
        raise IndexError(f'Not enough shapes in slide. is the slide type correct?')

def add_nstg_content_slide(prs,title_text = '', content_text = ''):
    cslide = prs.slides.add_slide(prs.slide_layouts[nstg_slide_layouts['content']])
    write_nstg_content_slide(cslide,title_text,content_text)
    return cslide  
    
def write_nstg_content_slide(slide,title_text = '', content_text = ''):
        
    try:
        slide.shapes[0].text = content_text
        slide.shapes[1].text = title_text
    except IndexError:
        raise IndexError(f'Not enough shapes in slide. is the slide type correct?')
        

def add_nstg_nocontent_slide(prs,title_text = ''):
    cslide = prs.slides.add_slide(prs.slide_layouts[nstg_slide_layouts['nocontent']])
    write_nstg_nocontent_slide(cslide,title_text)
    return cslide

def write_nstg_nocontent_slide(slide,title_text = ''):
        
    try:
        slide.shapes[0].text = title_text
    except IndexError:
        raise IndexError(f'Not enough shapes in slide. is the slide type correct?')

def add_tiled_ims(ppt,nocontentslide,imlist,title,gridshape, rowlabels = [], collabels = [], sep = 0, custmargins = {}):
    #imlist is paths, gridshape = (cols,rows)
    nims = len(imlist)
    
    for k in default_margins:
        if not k in custmargins:
            custmargins[k] = default_margins[k]
        else:
            if not isinstance(custmargins[k],Length):
                custmargins[k] = Inches(custmargins[k])
                
    if not isinstance(sep,Inches):
        sep = Inches(sep)
    
    collabels, rowlabels, fullgridshape = validate_labels(collabels,rowlabels,gridshape)
    
    # print(collabels,rowlabels)
    
    assert nims <= gridshape[0]*gridshape[1]
    
    pptdim = (ppt.slide_width,ppt.slide_height)
    
    workingsize = (pptdim[0]- custmargins['right'] - custmargins['left'],
                    pptdim[1]- custmargins['top'] - custmargins['bottom'])
    
    figsize = ((workingsize[0] - sep*(fullgridshape[0] - 1)) / fullgridshape[0],
               (workingsize[1] - sep*(fullgridshape[1] - 1)) / fullgridshape[1] )
    
    ##if labels, do labels
    txtfactor = 0.5
    coltbs = list()
    if collabels:
        
        txtsize = (figsize[0]*txtfactor,figsize[1]*txtfactor)
        for ind in range(len(collabels)):
            if not collabels[ind]:
                continue
            
            clabx = custmargins['left'] + (txtsize[0] + sep) + ind * (figsize[0] + sep) + sep//2 + (figsize[0] - txtsize[0])//2
            claby = custmargins['top'] + 0 * (figsize[1] + sep) + sep//2
            
            tb = add_text_box(nocontentslide,(clabx,claby),txtsize,collabels[ind])
            coltbs.append(tb)
            
        custmargins['top'] += txtsize[1]
    
    rowtbs= list()
    if rowlabels:
        
        txtsize = (figsize[0]*txtfactor,figsize[1]*txtfactor)
        for ind in range(len(rowlabels)):
            if not rowlabels[ind]:
                continue
            
            rlabx = custmargins['left'] + 0 * (figsize[0] + sep) + sep//2
            rlaby = custmargins['top'] + ind * (figsize[1] + sep) + sep//2 + (figsize[1] - txtsize[1])//2
        
            tb = add_text_box(nocontentslide,(rlabx,rlaby),txtsize,rowlabels[ind])
            
            rowtbs.append(tb)
        
        custmargins['left'] += txtsize[0]
    
    workingsize = (pptdim[0]- custmargins['right'] - custmargins['left'],
                    pptdim[1]- custmargins['top'] - custmargins['bottom'])
    
    figsize = ((workingsize[0] - sep*(gridshape[0] - 1)) / gridshape[0],
               (workingsize[1] - sep*(gridshape[1] - 1)) / gridshape[1] )
    
    ##adjust label position
    for i in range(len(coltbs)):
        coltbx = custmargins['left'] + i * (figsize[0] + sep) + sep//2 + figsize[0] // 2 - txtsize[0] //2
        coltbs[i].left = int(coltbx)
    
    for i in range(len(rowtbs)):
        rowtby = custmargins['top'] + i * (figsize[1] + sep) + sep//2 + figsize[1] // 2 - txtsize[1] //2
        rowtbs[i].top = int(rowtby)
    
    #adjust margins for new im grid
    #do ims
    for i in range(len(imlist)):
        ind = np.unravel_index(i,gridshape, order = 'F')
        
        imposx = custmargins['left'] + ind[0] * (figsize[0] + sep) + sep//2
        imposy = custmargins['top'] + ind[1] * (figsize[1] + sep) + sep//2
        
        if isinstance(imlist[i],str):
            imhandle = imlist[i]
            im = nocontentslide.shapes.add_picture(imhandle, top = int(imposy), left = int(imposx))
        
        elif isinstance(imlist[i],mpl.figure.Figure):
            imhandle = io.BytesIO()
            with imhandle as stream:
                imlist[i].savefig(stream, bbox_inches = 'tight',pad_inches = 0)
                im = nocontentslide.shapes.add_picture(stream, top = int(imposy), left = int(imposx))
            
        elif  isinstance(imlist[i],type(None)):
            size =  (figsize[0]//4,figsize[1]//4)
            size = (min(size),)*2
            pos = (int(imposx) +  figsize[0]//2 - size[0]//2,int(imposy) +  figsize[1]//2 - size[1]//2)
            noshape = nocontentslide.shapes.add_shape(MSO_SHAPE.NO_SYMBOL,left= pos[0], top = pos[1], width = size[0], height = size[1])
            fill = noshape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(128,0,0)
            
            line = noshape.line
            line.width = 0
            linefill= line.fill
            linefill.solid()
            linefill.fore_color.rgb = RGBColor(128,0,0)
            
            continue
        else:
            raise Exception(f'figure object not valid type: {type(imlist[i])}')
        
        imshape = (im.width,im.height)
        
        dim = np.argmax([imshape[0] / figsize[0], imshape[1] / figsize[1]])
        
        ratio = imshape[dim] / figsize[dim]
        imshape_new = (im.width / ratio,im.height / ratio)
        c = (-imshape_new[0] / 2 - sep / 2 + figsize[0] / 2, -imshape_new[1] / 2 - sep / 2 + figsize[1] / 2)

        im.width = int(im.width / ratio)
        im.height = int(im.height / ratio)
        im.left += int(c[0])
        im.top += int(c[1])
        
    nocontentslide.shapes[0].text = title
    return nocontentslide

def add_tiled_figs(ppt,nocontentslide,figlist,title,gridshape, sep = 0, custmargins = {}, savefigargs = {}):
    
    byteslist = list()
    for fig in figlist:
        
        f = io.BytesIO()
        with f as stream:
            fig.savefig(stream,**savefigargs)
        byteslist.append(f)
        
    return add_tiled_ims(ppt,nocontentslide,byteslist,title,gridshape, sep = sep, custmargins  = custmargins)

def set_line_properties(slide,shapetype = None, width = None, fillcolor = None, dashstyle = None):
    shapes = slide.shapes
    
    if not shapetype:
        shapetype = BaseShape
    
    for shape in shapes:
        if isinstance(shape,shapetype):
            line = shape.line
            line.width = Pt(width)
            line.color.rgb = RGBColor(*fillcolor)
            line.dash_style = dashstyle
        
    return slide

def add_text_box(slide,pos,size,text,**textargs):
    lbl = slide.shapes.add_textbox(left = pos[0], top = pos[1], width = size[0], height = size[1])
    
    lbl.text = text
    lbl.text_frame.vertical_anchor = 3
    for pg in lbl.text_frame.paragraphs:
        pg.alignment = 2
    lbl.text_frame.auto_size = 2
    
    # print(lbl.text)
    try:
        # lbl.text_frame.fit_text(max_size = 24)
        pass
    except TypeError:
        print(f'failed to format label {lbl.text}')
        # traceback.print_exc()
    # lbl.text_frame.fit_text()
    return lbl
    
    
def write_nstg_transition_slide(slide,trans_text = ''):        
    write_nstg_nocontent_slide(slide,title_text = trans_text)

#%%
def validate_labels(collabels,rowlabels, gridshape, reduce = False):
    
    collabels, gridshape = _validate_labels(collabels,gridshape,0, reduce = reduce)
    rowlabels, gridshape = _validate_labels(rowlabels,gridshape,1, reduce = reduce)
    
    return collabels,rowlabels,gridshape
    

def _validate_labels(labels,gridshape,ind, reduce = False):
    if labels:
        if len(labels) + 1 == gridshape[ind] and reduce:
            return labels,gridshape
        if len(labels) == gridshape[ind]:
            # gridshape = (gridshape[0], gridshape[1] + 1)
            pass
        elif len(labels) + 1 > gridshape[ind]:
            labels = labels[:gridshape[ind]]
        elif len(labels) < gridshape[ind]:
            diff = gridshape[ind] - len(labels)
            labels = labels + ['']*diff
        
        gridshape = tuple(gridshape[n] 
                          if not n == ind 
                          else gridshape[n] + 1 
                          for n in range(len(gridshape)))
        
        if not gridshape[ind] == len(labels) + 1:
            labels = []
        
        return labels, gridshape
    else:
        return [], gridshape